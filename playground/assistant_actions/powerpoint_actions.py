# This file contains the implementation of assistant actions related to PowerPoint presentations. It is still a work in progress.

from playground.actions_manager import agent_action
from playground.constants import ASSISTANTS_WORKING_FOLDER
from pptx import Presentation
from pptx.util import Inches
import markdown
import os


@agent_action
def create_new_powerpoint(filename):
    """
    Creates a new PowerPoint presentation and saves it with the specified filename.

    Parameters:
    filename (str): The name of the PowerPoint file to be created.

    Returns:
    None
    """
    template_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "template.pptx"
    )
    prs = Presentation(template_path)
    file_path = os.path.join(ASSISTANTS_WORKING_FOLDER, filename)
    prs.save(file_path)
    return f"Created new PowerPoint presentation: {filename}"


def set_placeholder_text(slide, placeholder_name, text):
    """
    Sets the text of a placeholder in a slide based on the placeholder name.

    Parameters:
    slide (Slide): The slide containing the placeholder.
    placeholder_name (str): The name of the placeholder to set the text for.
    text (str): The text to set in the placeholder.

    Returns:
    None
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.placeholder_format.idx == 0 and placeholder_name == "title":
            shape.text = text
            return
        elif shape.placeholder_format.idx == 1 and placeholder_name == "subtitle":
            shape.text = text
            return
        elif (
            placeholder_name == "content"
            and shape.text == ""
            and shape.name.startswith("Content")
        ):
            shape.text = text
    return


@agent_action
def add_slide_to_powerpoint(filename, slide_content):
    """
    Adds a new slide to an existing PowerPoint presentation.

    Parameters:
    filename (str): The name of the PowerPoint file to which the slide will be added.
    slide_content (dict): A dictionary containing the slide content with the following keys:
        - title (str): The title of the slide.
        - subtitle (str): The subtitle of the slide (if applicable).
        - content (list of str): A list of strings containing text content.
        - images (list of str): A list of image file paths to be added to the slide.
        - slide_type (str): The type of slide to add ("Title Slide", "Title and Content", "Lesson Summary", "Closing Slide").
        - notes (str): Notes to be added to the slide.

    Returns:
    None
    """
    file_path = os.path.join(ASSISTANTS_WORKING_FOLDER, filename)
    prs = Presentation(file_path)
    slide_type = slide_content.get("slide_type", "Title and Content")

    # Define slide layouts based on slide type
    slide_layouts = {
        "Title Slide": 0,  # Typically 0 is the layout for title slide
        "Title and Content": 3,  # Typically 1 is the layout for title and content
        "Lesson Summary": 8,  # Example layout index for lesson summary
        "Closing Slide": 14,  # Example layout index for closing slide
    }

    slide_layout_index = slide_layouts.get(slide_type, 1)
    slide_layout = prs.slide_layouts[slide_layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    if slide_content.get("title"):
        set_placeholder_text(slide, "title", slide_content.get("title"))

    # Set the subtitle if applicable
    if slide_type in ["Title Slide", "Lesson Summary", "Closing Slide"]:
        set_placeholder_text(slide, "subtitle", slide_content.get("subtitle", ""))

    # Add text content
    content_text = "\n".join(slide_content.get("content", []))
    html_content = markdown.markdown(content_text)
    set_placeholder_text(slide, "content", html_to_text(html_content))

    # Add images
    for image_path in slide_content.get("images", []):
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(1), Inches(2))

    # Add notes to the slide
    notes = slide_content.get("notes", "")
    slide.notes_slide.notes_text_frame.text = notes

    prs.save(file_path)
    return f'Added a slide of type "{slide_type}" to {filename}'


def html_to_text(html):
    """
    Converts HTML content to plain text.

    Parameters:
    html (str): The HTML content to be converted.

    Returns:
    str: The plain text extracted from the HTML content.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text()
